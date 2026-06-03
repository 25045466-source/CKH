#!/usr/bin/env python3
"""
Language Curriculum Design - Chapter 4: Principles
Editable PowerPoint Presentation (19 slides)
Based on Nation & Macalister (2010)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Initialize presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
COLOR_NAVY = RGBColor(30, 60, 114)       # Deep Navy (#1E3C72)
COLOR_BLUE = RGBColor(52, 152, 219)     # Muted Blue (#3498DB)
COLOR_GREY = RGBColor(248, 250, 252)    # Warm Grey (#F8FAFC)
COLOR_TEXT = RGBColor(44, 62, 80)       # Charcoal Text (#2C3E50)
COLOR_WHITE = RGBColor(255, 255, 255)   # White
COLOR_MUTED = RGBColor(127, 140, 141)   # Muted Grey

FONT_TITLE = "Arial"
FONT_BODY = "Calibri"

def add_solid_background(slide, color):
    """Add solid background to slide"""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg

def create_header(slide, title_text, is_dark=False):
    """Create header with title and decorative line"""
    tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.8))
    tf = tx_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = FONT_TITLE
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE if is_dark else COLOR_NAVY
    
    if not is_dark:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.4), Inches(11.733), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_BLUE
        line.line.fill.background()

def create_footer(slide, current_page, total_pages=19, is_dark=False):
    """Create footer with slide number"""
    tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.9), Inches(11.733), Inches(0.4))
    tf = tx_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Language Curriculum Design | Chapter 4: Principles                                                                                         Slide {current_page} / {total_pages}"
    p.font.name = FONT_BODY
    p.font.size = Pt(10)
    p.font.color.rgb = COLOR_BLUE if is_dark else COLOR_MUTED

# ============ SLIDE 1: Title Slide ============
blank_layout = prs.slide_layouts[6]
s1 = prs.slides.add_slide(blank_layout)
add_solid_background(s1, COLOR_NAVY)
t_box = s1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(2.5))
tf1 = t_box.text_frame
tf1.word_wrap = True
p1 = tf1.paragraphs[0]
p1.text = "A BLUEPRINT FOR CURRICULUM DESIGN"
p1.font.name = FONT_TITLE
p1.font.size = Pt(40)
p1.font.bold = True
p1.font.color.rgb = COLOR_WHITE

p2 = tf1.add_paragraph()
p2.text = "The Architecture of Language Learning: 20 Research-Backed Principles"
p2.font.name = FONT_BODY
p2.font.size = Pt(20)
p2.font.color.rgb = COLOR_BLUE
p2.space_before = Pt(15)

p3 = tf1.add_paragraph()
p3.text = "Based on Nation & Macalister (2010) - Chapter 4: Principles"
p3.font.name = FONT_BODY
p3.font.size = Pt(14)
p3.font.color.rgb = COLOR_WHITE
p3.space_before = Pt(20)
create_footer(s1, 1, is_dark=True)

# ============ SLIDE 2: Table of Contents ============
s2 = prs.slides.add_slide(blank_layout)
add_solid_background(s2, COLOR_GREY)
create_header(s2, "TABLE OF CONTENTS (AGENDA)")

agenda_items = [
    ("PART I: THEORETICAL FOUNDATIONS", 
     ["1. Beyond Methods: Moving to Research-Based Design", 
      "2. The Subdivision of Principles (Figure 4.1)", 
      "3. Method Approach vs. Empirical Evidence"]),
    ("PART II: THE THREE CORE SUBDIVISIONS", 
     ["4. Content and Sequencing Principles", 
      "5. Format and Presentation Principles", 
      "6. Monitoring and Assessment Principles"]),
    ("PART III: APPLICATIONS & SYNTHESIS", 
     ["7. Balancing the Four Strands", 
      "8. The Human Element in Curriculum", 
      "9. Executive Evaluation & References"])
]

for idx, (section, slides) in enumerate(agenda_items):
    col_x = Inches(0.8 + idx * 4.0)
    card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, col_x, Inches(1.8), Inches(3.7), Inches(4.8))
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_WHITE
    card.line.color.rgb = COLOR_BLUE
    card.line.width = Pt(1.5)
    
    tf = card.text_frame
    tf.word_wrap = True
    p_sec = tf.paragraphs[0]
    p_sec.text = section
    p_sec.font.name = FONT_TITLE
    p_sec.font.size = Pt(13)
    p_sec.font.bold = True
    p_sec.font.color.rgb = COLOR_NAVY
    p_sec.space_after = Pt(12)
    
    for s in slides:
        p_s = tf.add_paragraph()
        p_s.text = "• " + s
        p_s.font.name = FONT_BODY
        p_s.font.size = Pt(11)
        p_s.font.color.rgb = COLOR_TEXT
        p_s.space_after = Pt(8)
create_footer(s2, 2)

# ============ SLIDE 3: Section Divider - Part I ============
s3 = prs.slides.add_slide(blank_layout)
add_solid_background(s3, COLOR_NAVY)
t_box = s3.shapes.add_textbox(Inches(1.0), Inches(2.8), Inches(11.333), Inches(2.0))
tf = t_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "PART I: THEORETICAL FOUNDATIONS"
p.font.name = FONT_TITLE
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = COLOR_WHITE
p_sub = tf.add_paragraph()
p_sub.text = "Deconstructing Outdated Methods & Navigating Empirical Frameworks"
p_sub.font.name = FONT_BODY
p_sub.font.size = Pt(18)
p_sub.font.color.rgb = COLOR_BLUE
p_sub.space_before = Pt(10)
create_footer(s3, 3, is_dark=True)

# ============ SLIDE 4: Beyond Methods ============
s4 = prs.slides.add_slide(blank_layout)
add_solid_background(s4, COLOR_GREY)
create_header(s4, "BEYOND METHODS: RESEARCH-BASED DESIGN")
box = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.5))
box.fill.solid()
box.fill.fore_color.rgb = COLOR_WHITE
box.line.color.rgb = COLOR_BLUE
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "The Content Illusion & The Core Problem (Nation & Macalister, 2010)"
p.font.name = FONT_TITLE
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = COLOR_NAVY
p.space_after = Pt(14)

points = [
    "• The Fallacy of 'Modern' Courses: Many contemporary coursebooks repeat syllabi that vary only minimally from Berlitz models designed in the 1890s, failing to incorporate decades of empirical data regarding grammar item frequencies (George, 1963a, 1963b).",
    "• The Format Illusion: Innovative methods like Total Physical Response (TPR) and the Silent Way frequently introduce novel presentation formats while leaving underlying content selection unchanged (Richards & Rodgers, 1986).",
    "• Systematic Integrity: A truly sensible basis for curriculum avoids method-based distractions and instead translates peer-reviewed research directly into teaching guidelines (Nation & Macalister, 2010, p. 37)."
]
for pt in points:
    p_pt = tf.add_paragraph()
    p_pt.text = pt
    p_pt.font.name = FONT_BODY
    p_pt.font.size = Pt(13)
    p_pt.font.color.rgb = COLOR_TEXT
    p_pt.space_after = Pt(12)
create_footer(s4, 4)

# ============ SLIDE 5: The Subdivision of Principles ============
s5 = prs.slides.add_slide(blank_layout)
add_solid_background(s5, COLOR_GREY)
create_header(s5, "THE THREE SUBDIVISIONS OF PRINCIPLES")

subs = [
    ("CONTENT & SEQUENCING", "Concerned with what items are chosen for the language course, the order in which they are presented, and the depth of coverage. This group ensures that high-frequency, highly valuable items are prioritized to provide an optimal return on time invested (Nation & Macalister, 2010, p. 38)."),
    ("FORMAT & PRESENTATION", "Focused directly on classroom activities, pedagogical procedures, and how learners actively process instructional materials. This is the domain where classroom teachers exercise their greatest professional influence and autonomy."),
    ("MONITORING & ASSESSMENT", "Deals with tracking learner progress, executing diagnostic testing, and conducting comprehensive evaluation. This subdivision guarantees that learning is consistently scaffolded, feedback is targeted, and learning deficits are corrected.")
]

for idx, (sub_title, sub_desc) in enumerate(subs):
    x_pos = Inches(0.8 + idx * 4.0)
    card = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_pos, Inches(1.8), Inches(3.7), Inches(4.5))
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_WHITE
    card.line.color.rgb = COLOR_NAVY if idx == 0 else COLOR_BLUE
    tf = card.text_frame
    tf.word_wrap = True
    
    p_t = tf.paragraphs[0]
    p_t.text = sub_title
    p_t.font.name = FONT_TITLE
    p_t.font.size = Pt(14)
    p_t.font.bold = True
    p_t.font.color.rgb = COLOR_NAVY
    p_t.space_after = Pt(12)
    p_t.alignment = PP_ALIGN.CENTER
    
    p_d = tf.add_paragraph()
    p_d.text = sub_desc
    p_d.font.name = FONT_BODY
    p_d.font.size = Pt(12)
    p_d.font.color.rgb = COLOR_TEXT
create_footer(s5, 5)

# ============ SLIDE 6: Section Divider - Part II ============
s6 = prs.slides.add_slide(blank_layout)
add_solid_background(s6, COLOR_NAVY)
t_box = s6.shapes.add_textbox(Inches(1.0), Inches(2.8), Inches(11.333), Inches(2.0))
tf = t_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "PART II: THE THREE CORE SUBDIVISIONS"
p.font.name = FONT_TITLE
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = COLOR_WHITE
p_sub = tf.add_paragraph()
p_sub.text = "Deep-Dive into Content, Format, Presentation, and Assessment Principles"
p_sub.font.name = FONT_BODY
p_sub.font.size = Pt(18)
p_sub.font.color.rgb = COLOR_BLUE
p_sub.space_before = Pt(10)
create_footer(s6, 6, is_dark=True)

# ============ SLIDE 7: Content and Sequencing - Part 1 ============
s7 = prs.slides.add_slide(blank_layout)
add_solid_background(s7, COLOR_GREY)
create_header(s7, "CONTENT & SEQUENCING: PRINCIPLES 1-4")

principles_cs1 = [
    ("#1 Frequency & Utility", "A language course must explicitly prioritize high-frequency items (vocabulary, grammar, formulaic expressions) that yield the highest communicative value in regular language use (Nation & Macalister, 2010, p. 38; George, 1963a)."),
    ("#2 Sequencing Types", "Syllabus items can be sequenced using linear, cyclical, or modular approaches. Linear progression introduces items once; cyclical systems reintroduce items across expanding intervals to solidify retention."),
    ("#3 Interference Avoidance", "Curriculum must be structured to prevent the concurrent introduction of highly similar, confusing lexical or grammatical items (e.g., teaching antonyms or near-synonyms together), which can induce structural interference."),
    ("#4 Depth over Breadth", "Content must offer deep, repeated exposures to core items rather than superficial coverage of a vast quantity of language data, optimizing the cognitive return on instructional time invested.")
]

for idx, (p_t, p_d) in enumerate(principles_cs1):
    y_pos = Inches(1.8 + idx * 1.2)
    box = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y_pos, Inches(11.733), Inches(1.0))
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_WHITE
    box.line.color.rgb = COLOR_BLUE
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = p_t + " - "
    p.font.name = FONT_TITLE
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_NAVY
    run = p.add_run()
    run.text = p_d
    run.font.name = FONT_BODY
    run.font.size = Pt(11)
    run.font.color.rgb = COLOR_TEXT
create_footer(s7, 7)

# ============ SLIDE 8: Content and Sequencing - Part 2 ============
s8 = prs.slides.add_slide(blank_layout)
add_solid_background(s8, COLOR_GREY)
create_header(s8, "CONTENT & SEQUENCING: PRINCIPLES 5-8")

principles_cs2 = [
    ("#5 The Spaced Retrieval Matrix", "Items must be encountered across expanding intervals of time. Spaced retrieval forces learners to actively reconstruct cognitive paths to the target structure, boosting long-term memory allocation."),
    ("#6 Needs-Driven Adaptation", "Content must be selected and modified dynamically to match the objective socio-cognitive needs of specific learner sub-populations rather than adhering rigidly to standard textbook syllabi."),
    ("#7 Form-Meaning Connections", "Curriculum sequencing should cultivate strong links between grammatical forms and semantic meaning, ensuring that structures are never introduced as isolated grammatical entities."),
    ("#8 Theoretical Alignment", "The choice of what to teach must reflect contemporary findings in corpus linguistics, discourse analysis, and language acquisition research (Ellis, 2005; Nation & Macalister, 2010).")
]

for idx, (p_t, p_d) in enumerate(principles_cs2):
    y_pos = Inches(1.8 + idx * 1.2)
    box = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y_pos, Inches(11.733), Inches(1.0))
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_WHITE
    box.line.color.rgb = COLOR_BLUE
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = p_t + " - "
    p.font.name = FONT_TITLE
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_NAVY
    run = p.add_run()
    run.text = p_d
    run.font.name = FONT_BODY
    run.font.size = Pt(11)
    run.font.color.rgb = COLOR_TEXT
create_footer(s8, 8)

# ============ SLIDE 9: Format and Presentation - Four Strands ============
s9 = prs.slides.add_slide(blank_layout)
add_solid_background(s9, COLOR_GREY)
create_header(s9, "FORMAT & PRESENTATION: THE FOUR STRANDS")
box = s9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.5))
box.fill.solid()
box.fill.fore_color.rgb = COLOR_WHITE
box.line.color.rgb = COLOR_NAVY
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Balancing the 4 Strands of Language Learning (Nation & Macalister, 2010, p. 39)"
p.font.name = FONT_TITLE
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = COLOR_NAVY
p.space_after = Pt(14)

strands = [
    "1. Meaning-Focused Input: Learners must engage with rich oral and written language where their primary attention is focused on understanding messages (Comprehensible Input / i+1).",
    "2. Language-Focused Learning: Explicit attention is dedicated to the mechanics of the language system, including deliberate study of vocabulary, phonetics, and grammar rules.",
    "3. Meaning-Focused Output: Forcing production in speaking and writing across varied communicative contexts, ensuring learners produce messages that others can understand.",
    "4. Fluency Development: Helping learners maximize their use of structures they already know, working with familiar material under real-time processing pressures."
]
for st in strands:
    p_st = tf.add_paragraph()
    p_st.text = "• " + st
    p_st.font.name = FONT_BODY
    p_st.font.size = Pt(12)
    p_st.font.color.rgb = COLOR_TEXT
    p_st.space_after = Pt(10)
create_footer(s9, 9)

# ============ SLIDE 10: Format and Presentation - Part 2 ============
s10 = prs.slides.add_slide(blank_layout)
add_solid_background(s10, COLOR_GREY)
create_header(s10, "FORMAT & PRESENTATION: MOTIVATION & COGNITION")

items_fp = [
    ("Motivation & Engagement", "Ranked highly among formatting principles. Presentation must foster favorable attitudes toward the language, its users, the teacher's skill, and the learner's own chance of success (Integrative Motivation). Without interest and value, cognitive processing slows or halts entirely."),
    ("Time on Task & Focus", "Maximize the absolute quantity of time spent actively using the target language in the classroom. Avoid excessive use of the first language (L1) during meta-talk or administrative explanations."),
    ("Depth of Processing", "Encourage learners to process language items as deeply, thoughtfully, and analytically as possible, ensuring tasks require genuine semantic and cognitive engagement rather than mindless mechanical repetition.")
]

for idx, (t, d) in enumerate(items_fp):
    x_pos = Inches(0.8 + idx * 4.0)
    card = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(1.8), Inches(3.7), Inches(4.5))
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_WHITE
    card.line.color.rgb = COLOR_BLUE
    tf = card.text_frame
    tf.word_wrap = True
    p_t = tf.paragraphs[0]
    p_t.text = t
    p_t.font.name = FONT_TITLE
    p_t.font.size = Pt(14)
    p_t.font.bold = True
    p_t.font.color.rgb = COLOR_NAVY
    p_t.space_after = Pt(12)
    p_t.alignment = PP_ALIGN.CENTER
    
    p_d = tf.add_paragraph()
    p_d.text = d
    p_d.font.name = FONT_BODY
    p_d.font.size = Pt(12)
    p_d.font.color.rgb = COLOR_TEXT
create_footer(s10, 10)

# ============ SLIDE 11: Format and Presentation - Part 3 ============
s11 = prs.slides.add_slide(blank_layout)
add_solid_background(s11, COLOR_GREY)
create_header(s11, "FORMAT & PRESENTATION: PRINCIPLES 11-14")

principles_fp = [
    ("#11 Integrative Motivation", "Material presentation must deliberately align with learners' socio-integrative goals, helping them project positive future identities as competent speakers within the target language community."),
    ("#12 Learning Style Adaptability", "Provide diverse pedagogical paths and multi-sensory tasks to accommodate individual cognitive differences, learning rates, and sensory preferences (e.g., visual, auditory, kinesthetic processing)."),
    ("#13 Comprehensible Input (i+1)", "Ensure that approximately 95-98% of running words in input texts are already known by the learner, allowing them to decode the remaining unfamiliar items through context smoothly."),
    ("#14 Pushed Output Tasks", "Construct interaction patterns that compel learners to stretch their current linguistic resources, bridging gaps between receptive knowledge and productive capacity via collaborative output.")
]

for idx, (p_t, p_d) in enumerate(principles_fp):
    y_pos = Inches(1.8 + idx * 1.2)
    box = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y_pos, Inches(11.733), Inches(1.0))
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_WHITE
    box.line.color.rgb = COLOR_BLUE
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = p_t + " - "
    p.font.name = FONT_TITLE
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_NAVY
    run = p.add_run()
    run.text = p_d
    run.font.name = FONT_BODY
    run.font.size = Pt(11)
    run.font.color.rgb = COLOR_TEXT
create_footer(s11, 11)

# ============ SLIDE 12: Monitoring and Assessment - Part 1 ============
s12 = prs.slides.add_slide(blank_layout)
add_solid_background(s12, COLOR_GREY)
create_header(s12, "MONITORING & ASSESSMENT: THE CORE DRIVERS")

ma_drivers = [
    ("Diagnostic Tracking", "Continuous assessment should illuminate existing learner proficiencies and systemic gaps, guiding immediate pedagogical recalibrations rather than serving purely punitive grading purposes (Nation & Macalister, 2010, p. 40)."),
    ("Formative Feedback Loop", "Providing responsive, actionable feedback helps learners identify systematic errors and refine their internal language hypothesis mechanisms in real-time."),
    ("Washback Optimization", "Ensure that tests stimulate productive classroom behavior, beneficial study habits, and balanced communicative development rather than narrow, mechanical 'teaching to the test' practices.")
]

for idx, (t, d) in enumerate(ma_drivers):
    x_pos = Inches(0.8 + idx * 4.0)
    card = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(2.2), Inches(3.7), Inches(4.0))
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_WHITE
    card.line.color.rgb = COLOR_NAVY
    tf = card.text_frame
    tf.word_wrap = True
    p_t = tf.paragraphs[0]
    p_t.text = t
    p_t.font.name = FONT_TITLE
    p_t.font.size = Pt(14)
    p_t.font.bold = True
    p_t.font.color.rgb = COLOR_NAVY
    p_t.space_after = Pt(12)
    p_t.alignment = PP_ALIGN.CENTER
    
    p_d = tf.add_paragraph()
    p_d.text = d
    p_d.font.name = FONT_BODY
    p_d.font.size = Pt(12)
    p_d.font.color.rgb = COLOR_TEXT
create_footer(s12, 12)

# ============ SLIDE 13: Monitoring and Assessment - Part 2 ============
s13 = prs.slides.add_slide(blank_layout)
add_solid_background(s13, COLOR_GREY)
create_header(s13, "MONITORING & ASSESSMENT: PRINCIPLES 15-17")

principles_ma = [
    ("#15 Ongoing Formative Tracking", "Monitor learner progress dynamically using diverse, low-stakes tools (portfolios, logs, observation rubrics) to document real-world communicative growth over time."),
    ("#16 Criterion-Referenced Testing", "Evaluate performance against clear, explicit behavioral benchmarks and operational goals rather than ranking students against one another through norm-referenced curves."),
    ("#17 Empowering Self-Assessment", "Integrate metacognitive checkpoints to train learners in tracking their own progress, setting realistic micro-goals, and adjusting their autonomous study habits outside of class.")
]

for idx, (p_t, p_d) in enumerate(principles_ma):
    y_pos = Inches(2.0 + idx * 1.5)
    box = s13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y_pos, Inches(11.733), Inches(1.2))
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_WHITE
    box.line.color.rgb = COLOR_BLUE
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = p_t + " - "
    p.font.name = FONT_TITLE
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_NAVY
    run = p.add_run()
    run.text = p_d
    run.font.name = FONT_BODY
    run.font.size = Pt(12)
    run.font.color.rgb = COLOR_TEXT
create_footer(s13, 13)

# ============ SLIDE 14: Section Divider - Part III ============
s14 = prs.slides.add_slide(blank_layout)
add_solid_background(s14, COLOR_NAVY)
t_box = s14.shapes.add_textbox(Inches(1.0), Inches(2.8), Inches(11.333), Inches(2.0))
tf = t_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "PART III: APPLICATIONS & SYNTHESIS"
p.font.name = FONT_TITLE
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = COLOR_WHITE
p_sub = tf.add_paragraph()
p_sub.text = "Macro-Level Curriculum Synthesis and Pragmatic Classroom Recommendations"
p_sub.font.name = FONT_BODY
p_sub.font.size = Pt(18)
p_sub.font.color.rgb = COLOR_BLUE
p_sub.space_before = Pt(10)
create_footer(s14, 14, is_dark=True)

# ============ SLIDE 15: Cross-List Overlaps ============
s15 = prs.slides.add_slide(blank_layout)
add_solid_background(s15, COLOR_GREY)
create_header(s15, "CROSS-LIST OVERLAPS & CROSS-VALIDATION")

table_shape = s15.shapes.add_table(4, 3, Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.5))
table = table_shape.table
table.columns[0].width = Inches(3.0)
table.columns[1].width = Inches(4.366)
table.columns[2].width = Inches(4.366)

headers = ["Theoretical Concept", "Nation & Macalister (2010)", "Corroborating Frameworks (Ellis, Brown)"]
for col_idx, text in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.text = text
    cell.fill.solid()
    cell.fill.fore_color.rgb = COLOR_NAVY
    p = cell.text_frame.paragraphs[0]
    p.font.name = FONT_TITLE
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

matrix_data = [
    ["Input and Output Balance", "Principle #13 (Comprehensible Input) and Principle #14 (Pushed Output).", "Ellis (2005) - Principles of Instructed SLA regarding receptive and productive language balances."],
    ["Focus on Form", "Language-focused learning strand mixed evenly with meaning-focused strands.", "Long (1991), Ellis (2005) - Deliberate attention to linguistic form embedded naturally within communicative contexts."],
    ["Learner Autonomy", "Principle #17 (Self-Assessment) and dynamic learning style modifications.", "Brown (1993), Krahnke & Christison (1983) - Strategic investment and learner centeredness as core constructs."]
]

for r_idx, row in enumerate(matrix_data):
    for c_idx, text in enumerate(row):
        cell = table.cell(r_idx + 1, c_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_WHITE
        p = cell.text_frame.paragraphs[0]
        p.font.name = FONT_BODY
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_TEXT
create_footer(s15, 15)

# ============ SLIDE 16: Pragmatic Constraints ============
s16 = prs.slides.add_slide(blank_layout)
add_solid_background(s16, COLOR_GREY)
create_header(s16, "PRAGMATIC CONSTRAINTS IN IMPLEMENTATION")

constraints = [
    ("Institutional Mandates", "Standardized national curricula or rigid corporate benchmark exams can limit a teacher's flexibility to customize content and sequencing dynamically based on immediate learner needs."),
    ("Resource Availability", "Access to high-quality comprehensible input texts, digital corpus tools, or authentic listening materials varies greatly across global educational environments, restricting rich input strands."),
    ("Teacher Preparation & Training", "Successfully balancing the four strands and applying spaced retrieval strategies requires a deep, research-oriented understanding of applied linguistics that may exceed standard training models.")
]

for idx, (t, d) in enumerate(constraints):
    y_pos = Inches(1.8 + idx * 1.6)
    box = s16.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), y_pos, Inches(11.733), Inches(1.3))
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_WHITE
    box.line.color.rgb = COLOR_BLUE
    tf = box.text_frame
    tf.word_wrap = True
    p_t = tf.paragraphs[0]
    p_t.text = t
    p_t.font.name = FONT_TITLE
    p_t.font.size = Pt(13)
    p_t.font.bold = True
    p_t.font.color.rgb = COLOR_NAVY
    p_d = tf.add_paragraph()
    p_d.text = d
    p_d.font.name = FONT_BODY
    p_d.font.size = Pt(11)
    p_d.font.color.rgb = COLOR_TEXT
    p_d.space_before = Pt(4)
create_footer(s16, 16)

# ============ SLIDE 17: Strategic Recommendations ============
s17 = prs.slides.add_slide(blank_layout)
add_solid_background(s17, COLOR_GREY)
create_header(s17, "STRATEGIC RECOMMENDATIONS FOR DESIGNERS")
box = s17.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.5))
box.fill.solid()
box.fill.fore_color.rgb = COLOR_WHITE
box.line.color.rgb = COLOR_NAVY
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Actionable Steps for Modern Language Programs"
p.font.name = FONT_TITLE
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = COLOR_NAVY
p.space_after = Pt(12)

recs = [
    "1. Audit Existing Courses: Conduct empirical frequency counts on current syllabi vocabulary and grammar to ensure high-frequency coverage.",
    "2. Schedule Spaced Cycles: Intentionally build review modules into the calendar at expanding intervals (e.g., 1 week, 3 weeks, 6 weeks) to prevent structural decay.",
    "3. Equalize the Four Strands: Run a time-allocation check on daily lesson plans to ensure meaning-focused input and output are not eclipsed by grammar drilling.",
    "4. Diversify Assessment: Supplement summative end-of-term exams with criterion-referenced portfolio tracking to capture authentic student progress."
]
for r in recs:
    p_r = tf.add_paragraph()
    p_r.text = r
    p_r.font.name = FONT_BODY
    p_r.font.size = Pt(12)
    p_r.font.color.rgb = COLOR_TEXT
    p_r.space_after = Pt(10)
create_footer(s17, 17)

# ============ SLIDE 18: Executive Summary ============
s18 = prs.slides.add_slide(blank_layout)
add_solid_background(s18, COLOR_GREY)
create_header(s18, "EXECUTIVE SUMMARY: KEY TAKEAWAYS")

takes = [
    ("Empirical Over Methods", "Reject commercial method labels and instead build core syllabus architectures around peer-reviewed research findings in applied linguistics."),
    ("The Power of the 4 Strands", "Effective language learning requires balancing focus on form with comprehensive message decoding and active language production tasks."),
    ("Systemic Accountability", "Integrate continuous tracking and self-assessment mechanisms to align curriculum implementation with objective learner milestones.")
]

for idx, (t, d) in enumerate(takes):
    x_pos = Inches(0.8 + idx * 4.0)
    card = s18.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(2.0), Inches(3.7), Inches(4.3))
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_WHITE
    card.line.color.rgb = COLOR_BLUE
    tf = card.text_frame
    tf.word_wrap = True
    p_t = tf.paragraphs[0]
    p_t.text = t
    p_t.font.name = FONT_TITLE
    p_t.font.size = Pt(14)
    p_t.font.bold = True
    p_t.font.color.rgb = COLOR_NAVY
    p_t.space_after = Pt(12)
    p_t.alignment = PP_ALIGN.CENTER
    
    p_d = tf.add_paragraph()
    p_d.text = d
    p_d.font.name = FONT_BODY
    p_d.font.size = Pt(12)
    p_d.font.color.rgb = COLOR_TEXT
create_footer(s18, 18)

# ============ SLIDE 19: Comprehensive References ============
s19 = prs.slides.add_slide(blank_layout)
add_solid_background(s19, COLOR_GREY)
create_header(s19, "COMPREHENSIVE REFERENCES (APA 7th)")

ref_box = s19.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.8))
tf = ref_box.text_frame
tf.word_wrap = True

references_list = [
    "Brown, H. D. (1993). Principles of language learning and teaching (3rd ed.). Prentice Hall.",
    "Ellis, R. (2005). Principles of instructed language learning. Asian EFL Journal, 7(3), 9-24.",
    "George, H. V. (1963a). A verb-form frequency count. ELT Journal, 18(1), 31-37.",
    "George, H. V. (1963b). Monograph of the Institute of Education. Bandung.",
    "Krahnke, K. J., & Christison, M. A. (1983). Recent language research and some principles for language teachers. TESOL Quarterly, 17(4), 625-649.",
    "Long, M. H. (1991). Focus on form: A design feature in language teaching methodology. In K. de Bot, R. Ginsberg, & C. Kramsch (Eds.), Foreign language research in cross-cultural perspective (pp. 39-52). John Benjamins.",
    "Nation, I. S. P., & Macalister, J. (2010). Language curriculum design. Routledge.",
    "Richards, J. C., & Rodgers, T. S. (1986). Approaches and methods in language teaching. Cambridge University Press."
]

for idx, ref in enumerate(references_list):
    p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
    p.text = ref
    p.font.name = FONT_BODY
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_TEXT
    p.space_after = Pt(8)

create_footer(s19, 19)

# Save presentation
prs.save("Editable_C4_Presentation.pptx")
print("✅ SUCCESS: Editable_C4_Presentation.pptx created successfully!")
print("📥 File is ready for download and editing in PowerPoint, Google Slides, or LibreOffice Impress.")
